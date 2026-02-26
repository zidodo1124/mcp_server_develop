from typing import Optional, List, Dict
import os
import sys
import re
import math

# 确保项目根目录导入
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ========== 定制化停用词/黑名单（适配AI架构与系统PPT）==========
STOPWORDS = set(
    """
的 了 在 是 我 你 他 她 它 我们 你们 他们 这个 那个 这些 那些 和 与 或 及 而 但 也 
a an the and or but if is are was were be been being have has had do does did will would shall should may might must can could of in on at to for with by from as into through during before after above below up down out off over under again further then once here there when where why how all any both each few more most other some such no nor not only own same so than too very s t can will just don should now
的话 之类 等等 左右 上下 前后 多少 大小 高低 长短 强弱 好坏 与否 可以 能够 应该 必须 需要 想要 打算 准备 开始 结束 进行 完成 实现 达到 超过 不足 缺乏 拥有 具有 包含 包括 涉及 关于 对于 至于 由于 因为 所以 因此 从而 进而 于是 然后 接着 此外 另外 同时 同样 相反 否则 除非 只要 只有 如果 假如 假设 比如 例如 举例 比如 诸如 像 如 若 即 也就是 换句话说 总的来说 总而言之 综上所述
引论 计划 研究 成果 浙江大学 计算机学院 人工智能研究所 高等教育出版社 文档 内容 章节 部分 框架 体系 结构 方法 技术 过程 结果 问题 原因 影响 发展 应用 系统 开发 运行 设计 优化 计算 处理 训练 推理 加速 支持 提供 实现 构建 封装 分配 调度 管理 协同 工作 能力 性能 效率 能耗 带宽 延时 资源
John Hennessy David Patterson ACM Turing 作者 教授 院士 出版社 出版 年份 章节 页码 君子 年度 奖励 燃料 方法
"""
)

# 无意义前缀/后缀（AI领域定制）
MEANINGLESS_SUFFIX = {"的时代", "的发展", "的应用", "的技术", "的方法", "的过程", "的结果", "的问题", "的架构", "的框架"}
MEANINGLESS_PREFIX = {"基于", "通过", "利用", "采用", "实现", "达到", "超过", "不足", "缺乏", "面向", "针对"}

# AI架构领域专业名词白名单（强制保留，不会被误删）
AI_WHITELIST = {
    "GPU", "CPU", "TPU", "XPU", "DSP", "FLOPs", "SmartNIC", "人工智能芯片", "异构计算",
    "并行计算", "分布式训练", "数据并行", "张量并行", "流水线并行", "参数服务器",
    "Ring AllReduce", "类脑芯片", "智能交换机", "存算分离", "冯诺依曼结构", "人工智能算法",
    "人工智能架构与系统", "芯片", "框架", "微计算模式", "神经网络训练任务",
    "人工智能基础软硬件", "人工智能系统", "人工智能", "神经网络", "推理", "训练"
}

# 允许保留的词性（只保留名词类）
ALLOWED_POS = {
    "n",    # 名词
    "nr",   # 人名（如需关闭可注释）
    "ns",   # 地名
    "nt",   # 机构团体
    "nz",   # 其他专名
    "nw",   # 新词
    "vn",   # 名动词（如“训练”“推理”）
    "PER",  # 专有名词
    "LOC",
    "ORG"
}

# 核心实体黑名单（适配PPT，过滤冗余标题/通用词/人名/无意义短语）
ENTITY_BLACKLIST = {
    "介绍", "概述", "总结", "例子", "示例", "例如", "因此", "然后", "其中", "以及",
    "我们", "你们", "他们", "这个", "那个", "一种", "用于", "可以", "通过", "进行",
    "使用", "101计划", "人工智能引论", "君子性非异也", "善假于物", "延伸阅读", "咸与维新",
    "John", "Hennessy", "David", "Patterson", "ACM", "Turing", "作者", "教授", "院士",
    "数据是燃料", "模型是引擎", "算力是加速器", "和推理是数据密集", "年度图灵奖",
    "授予", "架构的", "以奖励其开创系统","奖励","君子","年度","燃料","方法"  # 新增过滤无意义短语
}

try:
    from tools import YA_MCPServer_Tool
except ImportError:
    def YA_MCPServer_Tool(**kwargs):
        def decorator(func):
            return func
        return decorator

try:
    from modules.YA_Common.utils.logger import get_logger
except ImportError:
    import logging
    def get_logger(name):
        logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(name)s - %(levelname)s - %(message)s")
        return logging.getLogger(name)

logger = get_logger("YA_MCPServer_KG_Tool-AI架构")

# ========== 配置项（核心调整：扩容边缘节点 + 极致减边）==========
def _load_kg_defaults() -> Dict:
    return {
        "top_k_entities": 35,          # 增加总节点数（为边缘节点预留空间）
        "min_occur": 2,                # 降低节点出现阈值，保留更多边缘节点
        "min_len": 2,
        "max_len": 10,
        "entity_blacklist": sorted(ENTITY_BLACKLIST),
        "min_edge_cooccurrence": 4,    # 提高共现阈值，大幅减少边
        "centrality_metric": "pagerank",
        "core_node_count": 5,          # 核心节点精简，突出重点
        "subcore_node_count": 15,      # 增加次核心节点数（边缘扩容1）
        "periphery_node_count": 10,    # 新增：独立外围节点数（边缘扩容2）
        "core_radius": 0.4,
        "subcore_radius": 1.5,
        "periphery_radius": 2.8,       # 外围层半径，物理隔离
        "edge_alpha": 0.12,            # 边更透明
        "edge_arc_rad": 0.02,
        "label_font_size": 10,
        # 节点大小：拉大差异，核心更突出
        "node_size_base": 600,
        "node_size_scale": 3000,
        "node_size_min": 400,          # 边缘节点最小尺寸
        "node_size_max": 4500,         # 核心节点最大尺寸
    }

# ========== 实体清洗与过滤（核心优化，解决关键字不准确/空节点）==========
def _clean_entity(entity: str) -> str:
    """AI领域定制实体清洗，移除无意义前缀/后缀，标准化缩写"""
    ent = entity.strip()
    # 移除无意义前缀
    for prefix in MEANINGLESS_PREFIX:
        if ent.startswith(prefix):
            ent = ent[len(prefix):].strip()
            break
    # 移除无意义后缀
    for suffix in MEANINGLESS_SUFFIX:
        if ent.endswith(suffix):
            ent = ent[:-len(suffix)].strip()
            break
    # AI领域缩写标准化（避免重复/错误）
    ent = ent.replace("GPU", "GPU").replace("TPU", "TPU").replace("XPU", "XPU")
    ent = ent.replace("AI芯片", "人工智能芯片").replace("异构计算", "异构计算")
    return ent

def _is_valid_entity(entity: str, min_len: int, max_len: int, blacklist: set) -> bool:
    """
    极致实体过滤规则（新增词性校验+白名单兜底）
    1. 强制保留AI专业名词白名单
    2. 过滤人名、机构名、无意义短语（如“xx是xx”）
    3. 过滤非名词类实体
    """
    # 基础清洗
    ent = _clean_entity(entity)
    
    # 强制保留白名单中的专业名词（最高优先级）
    if ent in AI_WHITELIST:
        return True
    
    # 空值/长度过滤
    if not ent or len(ent) < min_len or len(ent) > max_len:
        return False
    
    # 黑名单/停用词过滤（优先过滤无意义短语）
    if ent in STOPWORDS or ent in blacklist:
        return False
    
    # 过滤人名（包含英文大写开头的词）
    if re.match(r"^[A-Z][a-z]+$", ent):
        return False
    
    # 过滤“xx是xx”句式的无意义短语
    if re.search(r"是.+$", ent) and len(ent) > 8:
        return False
    
    # 纯符号/纯数字过滤（保留GPU/TPU等字母缩写）
    if re.fullmatch(r"[\W_]+", ent) or (re.fullmatch(r"\d+(\.\d+)?", ent) and not re.search(r"[A-Za-z]", ent)):
        return False
    
    # 过滤包含标点的碎片化文本
    if any(char in ent for char in "：；！？，。,.!?;:()（）[]【】"):
        return False
    
    # 过滤仅由虚词构成的短语
    function_words = {"的", "地", "得", "是", "了", "着", "过"}
    char_list = [c for c in ent if c in function_words]
    if len(char_list) >= len(ent) / 2:
        return False
    
    # 新增：词性校验（仅保留名词类）
    try:
        import jieba.posseg as pseg
        words = pseg.lcut(ent)
        for word, flag in words:
            # 只要有一个词不是允许的词性，就过滤
            if not any(flag.startswith(pos) for pos in ALLOWED_POS):
                return False
    except Exception:
        # jieba不可用时跳过词性校验，不影响基础过滤
        pass
    
    return True

# ========== 实体提取（新增词性标注+白名单强化）==========
def _extract_seed_entities_from_titles(slides: List[str], top_k_each: int = 5) -> List[str]:
    """从PPT标题提取种子实体，优先锁定核心概念"""
    seeds = []
    for s in slides:
        lines = [ln.strip() for ln in (s or "").splitlines() if ln and ln.strip()]
        if not lines:
            continue
        title = lines[0]
        # 过滤PPT冗余标题，只提取有效内容
        if any(kw in title for kw in ["人工智能引论", "第九章", "一、", "二、", "三、"]):
            title = re.sub(r"人工智能引论|第九章|[\u4e00-\u96f9]、", "", title).strip()
        if title:
            title_ents = _simple_entity_extraction(title, top_k=top_k_each)
            seeds.extend(title_ents)
    return list(dict.fromkeys(seeds))

def _simple_entity_extraction(text: str, top_k: int = 50) -> List[str]:
    """
    AI领域定制实体提取（核心优化：词性标注+白名单强化）
    1. 优先提取白名单中的专业名词
    2. 用jieba词性标注，只保留名词类
    3. 过滤通用动词/形容词/人名/无意义短语
    """
    if not text.strip():
        return []
    
    # 第一步：手动提取AI领域核心专业名词（兜底，确保关键概念不丢失）
    ai_pro_pattern = r"|".join([re.escape(term) for term in AI_WHITELIST])
    ai_pro_nouns = re.findall(ai_pro_pattern, text)
    
    # 第二步：jieba词性标注提取名词类实体
    try:
        import jieba
        import jieba.posseg as pseg
        # 分词并过滤词性
        words = pseg.cut(text)
        noun_entities = []
        for word, flag in words:
            # 只保留允许的词性
            if any(flag.startswith(pos) for pos in ALLOWED_POS):
                clean_word = _clean_entity(word)
                if (clean_word and 
                    clean_word not in STOPWORDS and 
                    clean_word not in ENTITY_BLACKLIST and
                    not re.match(r"^[A-Z][a-z]+$", clean_word) and  # 过滤人名
                    not re.search(r"是.+$", clean_word)):         # 过滤“xx是xx”
                    noun_entities.append(clean_word)
    except Exception as e:
        logger.warning(f"Jieba词性标注失败，使用旧版提取: {e}")
        # 降级方案：使用jieba analyse提取名词
        try:
            import jieba.analyse as analyse
            tags = analyse.extract_tags(
                text,
                topK=top_k + len(ai_pro_nouns),
                allowPOS=("n", "nr", "ns", "nt", "nz", "nw", "专名"),
                withWeight=False
            )
            noun_entities = [_clean_entity(t) for t in tags if t not in STOPWORDS]
        except Exception:
            noun_entities = []
    
    # 第三步：合并白名单和词性过滤结果，去重
    all_ents = ai_pro_nouns + noun_entities
    # 二次清洗和过滤
    all_ents = [_clean_entity(t) for t in all_ents if t]
    all_ents = [t for t in all_ents if _is_valid_entity(t, 2, 10, set())]
    # 去重并保留前top_k
    uniq_ents = list(dict.fromkeys(all_ents))
    return uniq_ents[:top_k]

# ========== 基础工具函数（保留，轻微优化）==========
def _compute_node_scores(graph, metric: str = "pagerank") -> Dict[str, float]:
    """计算节点重要性，优先PageRank，兜底度中心性"""
    try:
        import networkx as nx
    except Exception:
        return {n: 1.0 for n in graph.nodes()}
    if graph.number_of_nodes() == 0:
        return {}
    metric_l = (metric or "pagerank").lower()
    try:
        if metric_l == "degree":
            return {n: float(v) for n, v in graph.degree(weight="weight")}
        elif metric_l == "betweenness":
            return nx.betweenness_centrality(graph, normalized=True, weight="weight")
        else:
            return nx.pagerank(graph, alpha=0.85, weight="weight")
    except Exception:
        return {n: float(v) for n, v in graph.degree(weight="weight")}

def _extract_text_from_pptx(path: str) -> List[str]:
    """提取PPT文本，过滤空白幻灯片/冗余图片标注"""
    try:
        from pptx import Presentation
    except Exception:
        raise RuntimeError("python-pptx is required to parse PPTX files, 请安装：pip install python-pptx")
    if not os.path.exists(path):
        raise FileNotFoundError(f"PPT文件不存在：{path}")
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                # 过滤PPT图片标注/冗余字符
                if text and not re.match(r"tmp/|img|图\d+\.\d+", text):
                    parts.append(text)
        slide_text = "\n".join(parts)
        if slide_text:
            slides_text.append(slide_text)
        else:
            logger.warning("跳过空白/仅图片幻灯片")
    return slides_text

def _split_sentences_chinese(text: str) -> List[str]:
    """中文分句，适配PPT文本格式"""
    if not text:
        return []
    text = text.strip()
    parts = re.split(r"(?<=[。！？；\n])\s*", text)
    return [p.strip() for p in parts if p and len(p.strip()) > 5]

def _extract_key_sentences(text: str, top_k: int = 3) -> List[str]:
    """提取关键句，保留PPT核心内容"""
    if not text.strip():
        return []
    try:
        import jieba.analyse as analyse
    except Exception:
        return []
    try:
        tags = analyse.extract_tags(text, topK=20, withWeight=True)
    except Exception:
        tags = []
    kw_weights = {w[0]: w[1] for w in tags}
    sentences = _split_sentences_chinese(text)
    if not sentences:
        sentences = [s for s in text.splitlines() if s.strip()]
    scored = []
    for s in sentences:
        score = 0.0
        for kw, wt in kw_weights.items():
            if kw in s:
                score += wt
        scored.append((score, s))
    scored.sort(key=lambda x: x[0], reverse=True)
    selected = [s for _, s in scored if _ > 0][:top_k]
    if not selected:
        selected = sentences[:top_k]
    return selected

# ========== 核心函数1：抽取知识图谱（应用强共现过滤）==========
@YA_MCPServer_Tool(
    name="extract_knowledge_graph",
    title="Extract Knowledge Graph",
    description="从AI架构与系统PPT中抽取精准知识图谱，无空节点、关键字准确"
)
def extract_knowledge_graph(
    ppt_path: Optional[str] = None,
    text: Optional[str] = None,
    top_k_entities: int = 50,
    min_occur: Optional[int] = None,
    min_len: Optional[int] = None,
    max_len: Optional[int] = None,
    min_edge_cooccurrence: Optional[int] = None,  # 新增参数
    centrality_metric: Optional[str] = None,
    core_node_count: Optional[int] = None,
    subcore_node_count: Optional[int] = None,
    periphery_node_count: Optional[int] = None,   # 新增参数
    core_radius: Optional[float] = None,
    subcore_radius: Optional[float] = None,
    periphery_radius: Optional[float] = None,
    entity_blacklist: Optional[List[str]] = None,
) -> Dict:
    if not ppt_path and not text:
        raise ValueError("ppt_path 或 text 其中之一必须提供")
    # 加载配置
    defaults = _load_kg_defaults()
    top_k_entities = int(defaults.get("top_k_entities", top_k_entities))
    min_occur = int(defaults.get("min_occur", 2) if min_occur is None else min_occur)
    min_len = int(defaults.get("min_len", 2) if min_len is None else min_len)
    max_len = int(defaults.get("max_len", 10) if max_len is None else max_len)
    # 核心：强共现阈值
    min_edge_cooccurrence = int(defaults.get("min_edge_cooccurrence", 4) if min_edge_cooccurrence is None else min_edge_cooccurrence)
    centrality_metric = str(defaults.get("centrality_metric", "pagerank") if centrality_metric is None else centrality_metric)
    # 构建黑名单（基础+定制）
    blacklist = set(ENTITY_BLACKLIST)
    entity_blacklist_cfg = defaults.get("entity_blacklist", []) if entity_blacklist is None else entity_blacklist
    if isinstance(entity_blacklist_cfg, list):
        blacklist.update([str(i).strip() for i in entity_blacklist_cfg if str(i).strip()])
    # 提取PPT/文本内容
    slides = []
    if ppt_path:
        slides = _extract_text_from_pptx(ppt_path)
    else:
        slides = [text] if text and text.strip() else []
    slides = [s for s in slides if s.strip()]
    if not slides:
        logger.warning("无有效文本可提取")
        return {"nodes": [], "edges": [], "slides_summary": []}
    # 提取实体+种子实体（从标题提取，确保核心概念不丢失）
    slide_entities: List[List[str]] = []
    slide_summaries: List[Dict] = []
    all_entities = {}
    seed_entities = set(_extract_seed_entities_from_titles(slides))
    seed_entities = {e for e in seed_entities if _is_valid_entity(e, min_len, max_len, blacklist)}
    logger.info(f"从标题提取种子核心实体：{seed_entities}")
    # 遍历幻灯片提取实体
    for s in slides:
        raw_ents = _simple_entity_extraction(s, top_k=top_k_entities)
        # 严格过滤实体，确保无空节点/无意义节点
        ents = [e for e in raw_ents if _is_valid_entity(e, min_len, max_len, blacklist)]
        ents = list(dict.fromkeys(ents))  # 去重
        # 强制加入种子实体和白名单实体，确保核心概念保留
        ents = list(set(ents + list(seed_entities) + list(AI_WHITELIST)))
        slide_entities.append(ents)
        slide_summaries.append({"text": s, "key_points": _extract_key_sentences(s)})
        # 统计实体出现次数
        for e in ents:
            all_entities[e] = all_entities.get(e, 0) + 1
    # 过滤低频实体（降低阈值，保留更多边缘节点）
    filtered_entities = {e: cnt for e, cnt in all_entities.items() if cnt >= min_occur}
    if not filtered_entities:
        filtered_entities = dict(all_entities)
    ranked_entities = sorted(filtered_entities.items(), key=lambda x: x[1], reverse=True)
    ranked_entities = ranked_entities[: max(top_k_entities, 0)]
    candidate_set = {ent for ent, _ in ranked_entities}
    logger.info(f"过滤后候选实体（{len(candidate_set)}个）：{[e for e, _ in ranked_entities]}")
    # 构建边：极致去重 + 强共现过滤
    from collections import defaultdict
    edge_counts = defaultdict(lambda: defaultdict(int))
    for ents in slide_entities:
        uniq = [e for e in dict.fromkeys(ents) if e in candidate_set]
        for i in range(len(uniq)):
            for j in range(i + 1, len(uniq)):
                a, b = uniq[i], uniq[j]
                if a == b:
                    continue
                if a > b:  # 标准化边的key，避免(a,b)和(b,a)重复
                    a, b = b, a
                edge_counts[a][b] += 1
    # 核心：只保留共现次数≥阈值的边，彻底减少边数量
    pruned_edge_counts = {}
    for a, others in edge_counts.items():
        for b, w in others.items():
            if w >= min_edge_cooccurrence:
                pruned_edge_counts[(a, b)] = w

     # 1. 先找出所有孤立节点（没有边的节点）
    all_candidate_nodes = set(a for a, _ in ranked_entities)
    nodes_with_edges = set()
    for (a, b), _ in pruned_edge_counts.items():
        nodes_with_edges.add(a)
        nodes_with_edges.add(b)
    isolated_nodes = all_candidate_nodes - nodes_with_edges

    # 2. 为每个孤立节点，找到它共现次数最高的邻居，补一条边
    for node in isolated_nodes:
        # 找到该节点所有共现过的邻居
        neighbors = []
        if node in edge_counts:
            for neighbor, w in edge_counts[node].items():
                if neighbor != node:
                    neighbors.append((neighbor, w))
        # 也检查反向（node作为b的情况）
        for a, others in edge_counts.items():
            if node in others and a != node:
                neighbors.append((a, others[node]))
        if neighbors:
            # 选共现次数最高的邻居
            best_neighbor, best_w = max(neighbors, key=lambda x: x[1])
            # 标准化key
            if node > best_neighbor:
                key = (best_neighbor, node)
            else:
                key = (node, best_neighbor)
            # 补这条边（即使共现次数略低于阈值）
            if key not in pruned_edge_counts:
                pruned_edge_counts[key] = best_w

    # 3. 再按权重排序，保留前25条边（比之前的15条多，保证连通性）
    sorted_edges = sorted(pruned_edge_counts.items(), key=lambda x: x[1], reverse=True)
    pruned_edge_counts = dict(sorted_edges[:45])  # 从15条提高到25条，保证连通性

    # 计算节点重要性
    try:
        import networkx as nx
        candidate_graph = nx.Graph()
        for ent, cnt in ranked_entities:
            candidate_graph.add_node(ent, count=cnt)
        for (a, b), w in pruned_edge_counts.items():
            if a in candidate_graph and b in candidate_graph:
                candidate_graph.add_edge(a, b, weight=w)
        scores = _compute_node_scores(candidate_graph, metric=centrality_metric)
    except Exception as e:
        logger.warning(f"计算节点分数失败，使用词频兜底: {e}")
        scores = {ent: float(cnt) for ent, cnt in ranked_entities}
    # 构建最终节点和边
    sorted_by_score = sorted(ranked_entities, key=lambda x: (scores.get(x[0], 0.0), x[1]), reverse=True)
    final_ranked = sorted_by_score
    nodes = []
    id_map = {}
    for idx, (ent, cnt) in enumerate(final_ranked):
        nid = str(idx + 1)
        id_map[ent] = nid
        nodes.append({
            "id": nid,
            "label": ent,
            "count": cnt,  # 保留出现次数，用于节点大小缩放
            "score": round(float(scores.get(ent, 0.0)), 6),
        })
    edges = []
    for (a, b), w in pruned_edge_counts.items():
        if a in id_map and b in id_map:
            edges.append({
                "source": id_map[a],
                "target": id_map[b],
                "weight": w,
                "relation": "cooccurrence",
            })
    # 日志统计
    logger.info(
        f"图谱抽取完成：原始实体{len(all_entities)}个 → 过滤后有效节点{len(nodes)}个 | "
        f"原始边{sum(len(v) for v in edge_counts.values())}条 → 过滤后强关联边{len(edges)}条（共现阈值≥{min_edge_cooccurrence}）"
    )
    return {"nodes": nodes, "edges": edges, "slides_summary": slide_summaries}

# ========== 核心函数2：写入Neo4j（保留，兼容现有逻辑）==========
@YA_MCPServer_Tool(
    name="write_kg_to_neo4j",
    title="Write KG to Neo4j",
    description="将AI知识图谱写入Neo4j，merge节点与关系"
)
def write_kg_to_neo4j(
    kg: Dict,
    uri: str = "bolt://localhost:7687",
    user: str = "neo4j",
    password: str = "neo4j",
    node_label: str = "AI_Entity",
    rel_type: str = "RELATED_TO",
) -> Dict:
    if not kg or not kg.get("nodes") or not kg.get("edges"):
        logger.warning("空图谱，跳过写入Neo4j")
        return {"nodes_created": 0, "rels_created": 0}
    try:
        from py2neo import Graph, Node, Relationship
    except Exception:
        raise RuntimeError("py2neo 未安装，请安装：pip install py2neo")
    graph = Graph(uri, auth=(user, password))
    nodes = kg.get("nodes", [])
    edges = kg.get("edges", [])
    created_nodes = 0
    created_rels = 0
    # 写入节点（AI领域专属标签AI_Entity）
    for n in nodes:
        label = n.get("label")
        if not label:
            continue
        props = {"name": label, "count": n.get("count", 0), "score": n.get("score", 0.0), "domain": "人工智能架构与系统"}
        node = Node(node_label,** props)
        graph.merge(node, node_label, "name")
        created_nodes += 1
    # 写入关系（保留权重）
    for e in edges:
        src_id = e.get("source")
        tgt_id = e.get("target")
        try:
            src_idx = int(src_id) - 1
            tgt_idx = int(tgt_id) - 1
            src_label = nodes[src_idx]["label"]
            tgt_label = nodes[tgt_idx]["label"]
        except Exception:
            continue
        src_node = graph.nodes.match(node_label, name=src_label).first()
        tgt_node = graph.nodes.match(node_label, name=tgt_label).first()
        if src_node and tgt_node:
            rel = Relationship(
                src_node, rel_type, tgt_node,
                weight=e.get("weight", 1),
                relation=e.get("relation", "cooccurrence"),
                domain="人工智能架构与系统"
            )
            graph.merge(rel)
            created_rels += 1
    logger.info(f"Neo4j写入完成：节点{created_nodes}个，关系{created_rels}条")
    return {"nodes_created": created_nodes, "rels_created": created_rels}

# ========== 核心函数3：导出可视化（三层布局 + 动态节点大小）==========
@YA_MCPServer_Tool(
    name="export_kg_visualization",
    title="Export KG Visualization",
    description="三层布局+强关联边+动态节点大小，图谱清晰聚焦"
)
def export_kg_visualization(
    kg: Dict, path: str = "kg.graphml", format: str = "graphml"
) -> Dict:
    if not kg or not kg.get("nodes") or not kg.get("edges"):
        logger.warning("空图谱，跳过可视化导出")
        return {"path": path, "format": format, "error": None, "message": "空图谱，跳过导出"}
    try:
        import networkx as nx
    except Exception:
        raise RuntimeError("networkx 未安装，请安装：pip install networkx")
    # 创建输出目录
    out_dir = os.path.dirname(path)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)
        logger.info(f"创建导出目录：{out_dir}")

    # ========== 关键过滤：只保留有效节点和边 ==========
    nodes = kg.get("nodes", [])
    edges = kg.get("edges", [])

    # 1. 过滤节点：只保留有有效标签的实体
    valid_nodes = []
    valid_node_ids = set()
    node_label_map = {}  # id -> label
    node_count_map = {}  # label -> count（出现次数）
    for n in nodes:
        label = n.get("label")
        if label is None or not label.strip():
            continue
        # 额外过滤：纯符号/纯数字/图片标注/公式
        if re.fullmatch(r"[\W_]+", label) or re.fullmatch(r"\d+(\.\d+)?", label):
            continue
        if re.match(r"^图\d+.*$", label) or re.match(r"^(tmp|img).*$", label):
            continue
        if re.search(r"[∑∫√×÷±≠≈∞∂∇∈∉⊂⊃∪∩∅∀∃¬∧∨→↔⇒⇔∴∵∼≅≡≤≥]", label):
            continue
        valid_nodes.append(n)
        valid_node_ids.add(n["id"])
        node_label_map[n["id"]] = label
        node_count_map[label] = n.get("count", 1)  # 保存出现次数
    if not valid_nodes:
        logger.error("过滤后无有效节点，无法生成可视化")
        return {"path": path, "format": format, "error": "无有效节点"}

    # 2. 过滤边：只保留两端都是有效节点的强关联边
    valid_edges = []
    for e in edges:
        src_id = e.get("source")
        tgt_id = e.get("target")
        if src_id in valid_node_ids and tgt_id in valid_node_ids:
            valid_edges.append(e)

    # 3. 用过滤后的节点和边重建图
    G = nx.Graph()
    node_score_map = {}
    for n in valid_nodes:
        label = n["label"]
        score = float(n.get("score", 0.0))
        node_score_map[label] = score
        G.add_node(label, count=n.get("count", 1), score=score)
    for e in valid_edges:
        src_label = node_label_map[e["source"]]
        tgt_label = node_label_map[e["target"]]
        if src_label != tgt_label and src_label in node_score_map and tgt_label in node_score_map:
            G.add_edge(src_label, tgt_label, weight=e.get("weight", 1))

    # 导出逻辑
    fmt = format.lower()
    if fmt == "graphml":
        nx.write_graphml(G, path)
    elif fmt == "gexf":
        nx.write_gexf(G, path)
    elif fmt == "png":
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
            # 配置中文显示（强制生效，适配各种环境）
            plt.rcParams["font.sans-serif"] = ["SimHei", "Microsoft YaHei", "DejaVu Sans", "Arial Unicode MS"]
            plt.rcParams["axes.unicode_minus"] = False
            plt.rcParams['figure.dpi'] = 300
            plt.rcParams['savefig.dpi'] = 300
            plt.rcParams['savefig.facecolor'] = 'white'  # 白底，避免黑边
            plt.rcParams['savefig.bbox'] = 'tight'       # 紧凑布局，裁剪空白
        except Exception as e:
            raise RuntimeError(f"matplotlib 未安装/配置失败：{e}，请安装：pip install matplotlib")
        # 加载布局配置（三层布局参数）
        defaults = _load_kg_defaults()
        core_node_count = int(defaults.get("core_node_count", 5))
        subcore_node_count = int(defaults.get("subcore_node_count", 15))
        periphery_node_count = int(defaults.get("periphery_node_count", 10))
        core_radius = float(defaults.get("core_radius", 0.4))
        subcore_radius = float(defaults.get("subcore_radius", 1.5))
        periphery_radius = float(defaults.get("periphery_radius", 2.8))
        edge_alpha = float(defaults.get("edge_alpha", 0.12))
        edge_arc_rad = float(defaults.get("edge_arc_rad", 0.02))
        label_font_size = int(defaults.get("label_font_size", 10))
        # 节点大小缩放参数
        node_size_base = int(defaults.get("node_size_base", 600))
        node_size_scale = int(defaults.get("node_size_scale", 3000))
        node_size_min = int(defaults.get("node_size_min", 400))
        node_size_max = int(defaults.get("node_size_max", 4500))

        # ========== 核心：三层节点分层（核心+次核心+外围）==========
        sorted_nodes = sorted(node_score_map.items(), key=lambda x: x[1], reverse=True)
        # 1. 核心节点（最中心，最大）
        core_nodes = [n for n, _ in sorted_nodes[:core_node_count]]
        # 2. 次核心节点（中层，中等大小，数量增加）
        subcore_nodes = [n for n, _ in sorted_nodes[core_node_count:core_node_count+subcore_node_count]]
        # 3. 外围节点（外层，最小，新增独立层级，扩容边缘）
        periphery_nodes = [n for n, _ in sorted_nodes[core_node_count+subcore_node_count:core_node_count+subcore_node_count+periphery_node_count]]
        
        # 所有可见节点（三层合并）
        all_visible_nodes = core_nodes + subcore_nodes + periphery_nodes
        logger.info(f"三层布局：核心{len(core_nodes)}个 | 次核心{len(subcore_nodes)}个 | 外围{len(periphery_nodes)}个 | 总节点{len(all_visible_nodes)}个")

        # ========== 核心：分层布局计算（物理隔离，减少无效连接）==========
        pos = {}
        # 1. 核心节点：正中心紧凑布局
        core_count = len(core_nodes)
        for i, node in enumerate(core_nodes):
            angle = 2 * math.pi * i / core_count
            x = core_radius * math.cos(angle)
            y = core_radius * math.sin(angle)
            pos[node] = (x, y)
        # 2. 次核心节点：中层圆形布局
        subcore_count = len(subcore_nodes)
        for i, node in enumerate(subcore_nodes):
            angle = 2 * math.pi * i / subcore_count
            x = subcore_radius * math.cos(angle)
            y = subcore_radius * math.sin(angle)
            pos[node] = (x, y)
        # 3. 外围节点：外层圆形布局（新增，边缘扩容）
        periphery_count = len(periphery_nodes)
        for i, node in enumerate(periphery_nodes):
            angle = 2 * math.pi * i / periphery_count
            x = periphery_radius * math.cos(angle)
            y = periphery_radius * math.sin(angle)
            pos[node] = (x, y)

        # ========== 核心：二次过滤边（只保留可见节点的边）==========
        edges_to_draw = []
        for u, v, data in G.edges(data=True):
            if u in pos and v in pos:
                edges_to_draw.append((u, v, data))


        # 按权重排序，只保留前10条最强边（可视化层最终过滤）
        edges_to_draw.sort(key=lambda x: x[2].get("weight", 0), reverse=True)
        edges_to_draw = edges_to_draw[:100]  # 可视化只渲染前10条最强边

        # 重建可见子图
        G_visible = nx.Graph()
        for node in all_visible_nodes:
            if node in node_score_map:
                G_visible.add_node(node, score=node_score_map[node], count=node_count_map.get(node, 1))
        for u, v, data in edges_to_draw:
            G_visible.add_edge(u, v, weight=data.get("weight", 1))

        # ========== 节点样式：大小与出现次数强关联 + 三层颜色区分 ==========
        sizes = []
        node_colors = []
        node_edgecolors = []
        # 归一化出现次数（用于动态大小）
        max_count = max([d.get("count", 1) for _, d in G_visible.nodes(data=True)]) if G_visible.nodes() else 1
        
        for n, data in G_visible.nodes(data=True):
            count = data.get("count", 1)
            # 动态大小计算：出现次数越多，节点越大
            if max_count > 1:
                norm_count = (count - 1) / (max_count - 1)
            else:
                norm_count = 0.5
            node_size = node_size_base + node_size_scale * norm_count
            node_size = max(node_size_min, min(node_size_max, node_size))
            sizes.append(node_size)
            
            # 三层颜色区分
            if n in core_nodes:
                node_colors.append("#1f77b4")  # 深蓝（核心）
                node_edgecolors.append("#ffffff")
            elif n in subcore_nodes:
                node_colors.append("#ff7f0e")  # 橙色（次核心）
                node_edgecolors.append("#ffffff")
            else:
                node_colors.append("#2ca02c")  # 绿色（外围/边缘）
                node_edgecolors.append("#ffffff")

        # ========== 边样式：只保留强关联，粗细按权重 ==========
        edge_weights = [float(d.get("weight", 1)) for u, v, d in G_visible.edges(data=True)]
        max_w = max(edge_weights) if edge_weights else 1.0
        edge_widths = [0.5 + 2.5 * (w / max_w) for w in edge_weights]  # 强关联边更粗，弱关联更细

        # ========== 标签配置：核心全显，边缘按需显示 ==========
        labels = {}
        # 核心节点必显
        for n in core_nodes:
            labels[n] = n
        # 次核心节点必显
        for n in subcore_nodes:
            labels[n] = n
        # 外围节点可选（避免文字拥挤，此处选择显示）
        for n in periphery_nodes:
            labels[n] = n

        # ========== 绘图执行 ==========
        plt.figure(figsize=(22, 18))
        # 绘制边：弧形+低透明，仅显示强关联
        nx.draw_networkx_edges(
            G_visible, pos,
            alpha=edge_alpha,
            edge_color="#666666",
            width=edge_widths,
            connectionstyle=f"arc3,rad={edge_arc_rad}"
        )
        # 绘制节点：动态大小+三层颜色
        nx.draw_networkx_nodes(
            G_visible, pos,
            node_size=sizes,
            node_color=node_colors,
            alpha=0.9,
            edgecolors=node_edgecolors,
            linewidths=2
        )
        # 绘制标签：清晰背景
        nx.draw_networkx_labels(
            G_visible, pos, labels,
            font_size=label_font_size,
            font_family="sans-serif",
            font_weight="bold",
            font_color="white",
            bbox=dict(boxstyle="round,pad=0.25", facecolor="black", alpha=0.7)
        )
        # 美化
        plt.axis("off")
        plt.tight_layout()
        plt.savefig(path, dpi=300, bbox_inches="tight", pad_inches=0.1)
        plt.close()
        logger.info(f"PNG可视化导出完成：{path}")
    else:
        raise ValueError(f"不支持的导出格式: {format}，仅支持graphml/gexf/png")
    return {"path": path, "format": fmt, "node_count": len(G.nodes()), "edge_count": len(G.edges())}

# ========== 核心函数4：一站式处理（传递所有参数）==========
@YA_MCPServer_Tool(
    name="process_and_publish_kg",
    title="Process and Publish KG",
    description="一站式处理AI架构PPT：抽取图谱+导出可视化+可选写入Neo4j"
)
def process_and_publish_kg(
    ppt_path: Optional[str] = None,
    text: Optional[str] = None,
    top_k_entities: int = 50,
    min_occur: Optional[int] = None,
    min_len: Optional[int] = None,
    max_len: Optional[int] = None,
    min_edge_cooccurrence: Optional[int] = None,
    centrality_metric: Optional[str] = None,
    core_node_count: Optional[int] = None,
    subcore_node_count: Optional[int] = None,
    periphery_node_count: Optional[int] = None,
    core_radius: Optional[float] = None,
    subcore_radius: Optional[float] = None,
    periphery_radius: Optional[float] = None,
    entity_blacklist: Optional[List[str]] = None,
    write_neo4j: bool = False,
    neo_uri: str = "bolt://localhost:7687",
    neo_user: str = "neo4j",
    neo_password: str = "neo4j",
    export_format: str = "graphml",
    export_path: str = "ai_kg_output",
) -> Dict:
    logger.info(f"开始一站式处理AI架构与系统图谱 | PPT路径：{ppt_path} | 导出格式：{export_format}")
    neo_result = None
    export_result = None
    try:
        # 1. 抽取图谱（传递强共现阈值）
        kg = extract_knowledge_graph(
            ppt_path=ppt_path,
            text=text,
            top_k_entities=top_k_entities,
            min_occur=min_occur,
            min_len=min_len,
            max_len=max_len,
            min_edge_cooccurrence=min_edge_cooccurrence,
            centrality_metric=centrality_metric,
            core_node_count=core_node_count,
            subcore_node_count=subcore_node_count,
            periphery_node_count=periphery_node_count,
            core_radius=core_radius,
            subcore_radius=subcore_radius,
            periphery_radius=periphery_radius,
            entity_blacklist=entity_blacklist,
        )
        # 2. 可选写入Neo4j
        if write_neo4j:
            logger.info("开始写入Neo4j图数据库")
            try:
                neo_result = write_kg_to_neo4j(kg, uri=neo_uri, user=neo_user, password=neo_password)
            except Exception as e:
                neo_result = {"error": str(e)}
                logger.error(f"Neo4j写入失败: {e}", exc_info=True)
        # 3. 导出可视化
        logger.info("开始导出图谱可视化")
        fmt = export_format.lower()
        # 拼接导出路径（自动加后缀）
        if fmt == "png":
            out_file = f"{export_path}.png"
        elif fmt == "gexf":
            out_file = f"{export_path}.gexf"
        else:
            out_file = f"{export_path}.graphml"
        try:
            export_result = export_kg_visualization(kg, path=out_file, format=fmt)
        except Exception as e:
            export_result = {"error": str(e)}
            logger.error(f"可视化导出失败: {e}", exc_info=True)
        # 处理完成
        logger.info("AI架构与系统图谱一站式处理完成！")
        return {
            "kg": kg,
            "neo_result": neo_result,
            "export_result": export_result,
            "status": "success" if not (neo_result and "error" in neo_result) and not (export_result and "error" in export_result) else "partial_success"
        }
    except Exception as exc:
        logger.error(f"一站式处理失败: {exc}", exc_info=True)
        return {
            "kg": {"nodes": [], "edges": [], "slides_summary": []},
            "neo_result": None,
            "export_result": {"error": str(exc)},
            "status": "failed"
        }

# ========== 测试入口 ==========
if __name__ == "__main__":
    # 测试：替换为你的PPT路径
    TEST_PPT_PATH = "1.pptx"
    # 一站式处理：强过滤+三层布局
    result = process_and_publish_kg(
        ppt_path=TEST_PPT_PATH,
        export_format="png",
        export_path="ai_architecture_kg_final",
        min_edge_cooccurrence=4,  # 强共现阈值，边极少
        write_neo4j=False
    )
    # 打印结果统计
    print("="*60)
    print("AI架构与系统知识图谱（最终优化版）提取结果")
    print("="*60)
    print(f"处理状态：{result['status']}")
    print(f"有效节点数：{len(result['kg']['nodes'])}")
    print(f"强关联边数：{len(result['kg']['edges'])}")
    print(f"可视化文件：{result['export_result'].get('path', '未导出')}")
    if result['kg']['nodes']:
        print(f"核心节点（按重要性）：{[n['label'] for n in result['kg']['nodes'][:5]]}")