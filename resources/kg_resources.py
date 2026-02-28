# resources/kg_resources.py
import os
from typing import Any, Optional
from resources import YA_MCPServer_Resource


# 配置PPT文件存放目录
PPT_ROOT = "data/ppt/"  # 存放PPT文件的目录
KG_OUTPUT_ROOT = "output/kg/"  # 存放知识图谱输出结果的目录


@YA_MCPServer_Resource(
    "file:///ppt/{filename}",  # 资源模板 URI
    name="get_ppt_file",
    title="Get PPT File",
    description="获取指定路径的PPT文件内容（文本提取）",
)
def get_ppt_file(filename: str) -> Any:
    """
    返回PPT文件的文本内容（用于查看PPT中的文字）。

    Args:
        filename (str): PPT文件名，如 "ai_chapter1.pptx"

    Returns:
        Any: PPT文件的文本内容或错误信息。
    """
    try:
        from pptx import Presentation
        
        # 安全性检查：防止路径遍历攻击
        safe_filename = os.path.basename(filename)  # 只取文件名，忽略路径
        full_path = os.path.normpath(os.path.join(PPT_ROOT, safe_filename))
        
        # 确保文件在PPT_ROOT目录下
        if not full_path.startswith(os.path.abspath(PPT_ROOT)):
            return {"error": "Invalid file path"}
        
        if not os.path.exists(full_path):
            return {"error": f"PPT file not found: {safe_filename}"}
        
        # 提取PPT文本内容
        prs = Presentation(full_path)
        slides_text = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            if slide_content:
                slides_text.append({
                    "slide_number": i + 1,
                    "content": "\n".join(slide_content)
                })
        
        return {
            "filename": safe_filename,
            "slide_count": len(slides_text),
            "slides": slides_text
        }
    except ImportError:
        return {"error": "python-pptx is required, please install: pip install python-pptx"}
    except Exception as e:
        return {"error": str(e)}


@YA_MCPServer_Resource(
    "file:///ppt/list",  # 资源 URI
    name="list_ppt_files",
    title="List PPT Files",
    description="列出所有可用的PPT文件",
)
def list_ppt_files() -> Any:
    """
    列出PPT目录中的所有PPT文件。

    Returns:
        Any: PPT文件列表。
    """
    try:
        if not os.path.exists(PPT_ROOT):
            os.makedirs(PPT_ROOT, exist_ok=True)
            return {"files": [], "message": "PPT directory created, please add PPT files"}
        
        files = []
        for f in os.listdir(PPT_ROOT):
            if f.endswith(('.pptx', '.ppt')):
                file_path = os.path.join(PPT_ROOT, f)
                files.append({
                    "name": f,
                    "size": os.path.getsize(file_path),
                    "modified": os.path.getmtime(file_path)
                })
        
        return {
            "directory": PPT_ROOT,
            "file_count": len(files),
            "files": files
        }
    except Exception as e:
        return {"error": str(e)}


@YA_MCPServer_Resource(
    "file:///kg/output/{filename}",  # 资源模板 URI
    name="get_kg_output",
    title="Get KG Output",
    description="获取生成的知识图谱输出文件",
)
def get_kg_output(filename: str) -> Any:
    """
    返回知识图谱输出文件内容。

    Args:
        filename (str): 输出文件名，如 "ai_kg_output.png" 或 "ai_kg_output.graphml"

    Returns:
        Any: 文件内容或错误信息。
    """
    try:
        # 安全性检查
        safe_filename = os.path.basename(filename)
        full_path = os.path.normpath(os.path.join(KG_OUTPUT_ROOT, safe_filename))
        
        # 确保文件在KG_OUTPUT_ROOT目录下
        if not full_path.startswith(os.path.abspath(KG_OUTPUT_ROOT)):
            return {"error": "Invalid file path"}
        
        if not os.path.exists(full_path):
            return {"error": f"File not found: {safe_filename}"}
        
        # 根据文件类型返回不同内容
        if safe_filename.endswith(('.png', '.jpg', '.jpeg', '.gif')):
            # 对于图片文件，返回二进制数据（需要特殊处理）
            with open(full_path, 'rb') as f:
                import base64
                image_data = base64.b64encode(f.read()).decode('utf-8')
                return {
                    "filename": safe_filename,
                    "format": "image",
                    "data": image_data,
                    "size": os.path.getsize(full_path)
                }
        else:
            # 对于文本文件，直接返回内容
            with open(full_path, 'r', encoding='utf-8') as f:
                content = f.read()
                return {
                    "filename": safe_filename,
                    "format": "text",
                    "content": content,
                    "size": os.path.getsize(full_path)
                }
    except Exception as e:
        return {"error": str(e)}


@YA_MCPServer_Resource(
    "file:///kg/list",  # 资源 URI
    name="list_kg_outputs",
    title="List KG Outputs",
    description="列出所有生成的知识图谱输出文件",
)
def list_kg_outputs() -> Any:
    """
    列出知识图谱输出目录中的所有文件。

    Returns:
        Any: 输出文件列表。
    """
    try:
        if not os.path.exists(KG_OUTPUT_ROOT):
            os.makedirs(KG_OUTPUT_ROOT, exist_ok=True)
            return {"files": [], "message": "Output directory created"}
        
        files = []
        for f in os.listdir(KG_OUTPUT_ROOT):
            file_path = os.path.join(KG_OUTPUT_ROOT, f)
            files.append({
                "name": f,
                "size": os.path.getsize(file_path),
                "modified": os.path.getmtime(file_path),
                "type": os.path.splitext(f)[1][1:]  # 文件扩展名
            })
        
        return {
            "directory": KG_OUTPUT_ROOT,
            "file_count": len(files),
            "files": files
        }
    except Exception as e:
        return {"error": str(e)}


@YA_MCPServer_Resource(
    "file:///kg/preview/{filename}",  # 资源模板 URI
    name="preview_kg_graph",
    title="Preview KG Graph",
    description="预览知识图谱（返回PNG图片的base64编码）",
)
def preview_kg_graph(filename: str) -> Any:
    """
    预览知识图谱图片。

    Args:
        filename (str): PNG文件名，如 "ai_kg_output.png"

    Returns:
        Any: 图片的base64编码或错误信息。
    """
    try:
        # 只允许预览PNG文件
        if not filename.endswith('.png'):
            return {"error": "Only PNG files can be previewed"}
        
        safe_filename = os.path.basename(filename)
        full_path = os.path.normpath(os.path.join(KG_OUTPUT_ROOT, safe_filename))
        
        if not full_path.startswith(os.path.abspath(KG_OUTPUT_ROOT)):
            return {"error": "Invalid file path"}
        
        if not os.path.exists(full_path):
            return {"error": f"File not found: {safe_filename}"}
        
        # 读取图片并转换为base64
        with open(full_path, 'rb') as f:
            import base64
            image_data = base64.b64encode(f.read()).decode('utf-8')
            return {
                "filename": safe_filename,
                "format": "png",
                "data": image_data,
                "size": os.path.getsize(full_path)
            }
    except Exception as e:
        return {"error": str(e)}